"""Extract text from uploaded materials (PDF, PPTX, video transcription)."""

import io
import tempfile
import subprocess
from typing import Optional

from PyPDF2 import PdfReader
from pptx import Presentation

from app.models.content import UploadedMaterial
from app.services.storage import get_file_from_s3


async def extract_text_from_material(material: UploadedMaterial) -> Optional[str]:
    if material.extracted_text:
        return material.extracted_text

    try:
        file_bytes = await get_file_from_s3(material.s3_key)
    except Exception:
        return None

    if material.file_type == "pdf":
        return _extract_pdf(file_bytes)
    elif material.file_type == "pptx":
        return _extract_pptx(file_bytes)
    elif material.file_type == "mp4":
        return await _transcribe_video(file_bytes)
    return None


def _extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)


def _extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


async def _transcribe_video(data: bytes) -> str:
    """Extract audio from video and transcribe.
    For a production pipeline, this would use a real ASR service
    (e.g., MiniMax speech-to-text or Whisper). Currently returns a
    placeholder indicating that transcription infrastructure is needed."""
    # TODO: Integrate real ASR (Whisper / MiniMax speech-to-text)
    # Steps for real implementation:
    # 1. Save video to temp file
    # 2. Extract audio with ffmpeg: ffmpeg -i input.mp4 -vn -acodec pcm_s16le output.wav
    # 3. Send audio to ASR API
    # 4. Return transcript text
    return "[Video transcription placeholder — integrate ASR service for real transcripts]"
